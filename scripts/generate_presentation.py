"""
IntelliFit Graduation Project - PowerPoint Presentation Generator

Generates a professional academic graduation project presentation
following the standard university presentation structure.

Usage:
    python scripts/generate_presentation.py [output_path]

Customize team info by editing the PROJECT_INFO dictionary below.

Dependencies:
    pip install python-pptx
"""

import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ── Project Info (customize these) ───────────────────────────────────────────
PROJECT_INFO = {
    "team_members": [
        "Youssef Essam",
        # Add remaining team member names below:
        "Team Member 2",
        "Team Member 3",
        "Team Member 4",
    ],
    "supervisor": "Dr. [Supervisor Name]",
    "university": "[University Name]",
    "faculty": "Faculty of Computer Science",
    "year": "2024 – 2025",
    "github_url": "github.com/YoussefEssam74/Graduation-project",
}


# ── Color Palette ────────────────────────────────────────────────────────────
PRIMARY = RGBColor(0x1A, 0x1A, 0x2E)       # Dark navy
ACCENT = RGBColor(0x00, 0x7B, 0xFF)        # Bright blue
ACCENT_DARK = RGBColor(0x00, 0x5A, 0xCC)   # Darker blue
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF5)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
SUBTITLE_GRAY = RGBColor(0x66, 0x66, 0x66)
SUCCESS_GREEN = RGBColor(0x28, 0xA7, 0x45)


# ── Helper Functions ─────────────────────────────────────────────────────────

def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, transparency=0):
    """Add a colored rectangle shape to a slide."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_title_bar(slide, title_text):
    """Add a styled title bar at the top of content slides."""
    bar = add_shape_bg(slide, Inches(0), Inches(0), Inches(10), Inches(1.1), PRIMARY)
    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    accent_bar = add_shape_bg(
        slide, Inches(0), Inches(1.1), Inches(10), Inches(0.05), ACCENT
    )
    return bar


def add_bullet_content(slide, left, top, width, height, items, font_size=16,
                       color=DARK_TEXT, bold_first=False, line_spacing=1.5):
    """Add a text box with bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
        p.line_spacing = Pt(font_size * line_spacing)
        if bold_first and i == 0:
            p.font.bold = True
    return txBox


def add_section_header(slide, left, top, width, text, font_size=20):
    """Add a section header text."""
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = ACCENT_DARK
    return txBox


def add_footer(slide, slide_number, total_slides):
    """Add a footer with slide number."""
    footer_bar = add_shape_bg(
        slide, Inches(0), Inches(7.1), Inches(10), Inches(0.4), PRIMARY
    )
    txBox = slide.shapes.add_textbox(Inches(8.5), Inches(7.15), Inches(1.3), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{slide_number} / {total_slides}"
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.RIGHT

    txBox2 = slide.shapes.add_textbox(Inches(0.3), Inches(7.15), Inches(3), Inches(0.3))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "IntelliFit – Smart Gym Management System"
    p2.font.size = Pt(9)
    p2.font.color.rgb = WHITE
    p2.font.italic = True


# ── Slide Builders ───────────────────────────────────────────────────────────

TOTAL_SLIDES = 22


def slide_01_title(prs):
    """Slide 1: Project Title Slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_bg(slide, PRIMARY)

    # Accent line
    add_shape_bg(slide, Inches(1), Inches(1.8), Inches(2), Inches(0.05), ACCENT)

    # Project name
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "IntelliFit"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE

    p2 = tf.add_paragraph()
    p2.text = "Smart Gym Management System with AI-Powered Coaching"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(0xAA, 0xCC, 0xFF)

    # Team members
    info_items = ["Team Members:"]
    for member in PROJECT_INFO["team_members"]:
        info_items.append(f"    • {member}")
    info_items.extend([
        "",
        f"Supervisor: {PROJECT_INFO['supervisor']}",
        f"University: {PROJECT_INFO['university']} – {PROJECT_INFO['faculty']}",
        f"Year: {PROJECT_INFO['year']}",
    ])
    add_bullet_content(slide, Inches(1), Inches(3.5), Inches(8), Inches(3.5),
                       info_items, font_size=14, color=RGBColor(0xCC, 0xCC, 0xDD))


def slide_02_toc(prs):
    """Slide 2: Table of Contents."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Table of Contents")

    left_items = [
        "1.  Introduction",
        "2.  Problem Definition",
        "3.  Motivation",
        "4.  Objectives",
        "5.  Technology Comparison",
        "6.  System Lifecycle / Workflow",
        "7.  Similar Applications",
        "8.  Methodology",
        "9.  Timeline",
        "10. Functional Requirements",
        "11. System Design",
    ]
    right_items = [
        "12. Technologies Used",
        "13. Database Design",
        "14. UI / Screens",
        "15. Project Pipeline",
        "16. Progress Percentage",
        "17. Conclusion",
        "18. Challenges",
        "19. Future Work",
        "20. Thank You",
    ]
    add_bullet_content(slide, Inches(0.5), Inches(1.5), Inches(4.5), Inches(5),
                       left_items, font_size=15, color=DARK_TEXT, line_spacing=1.6)
    add_bullet_content(slide, Inches(5.2), Inches(1.5), Inches(4.5), Inches(5),
                       right_items, font_size=15, color=DARK_TEXT, line_spacing=1.6)
    add_footer(slide, 2, TOTAL_SLIDES)


def slide_03_introduction(prs):
    """Slide 3: Introduction."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Introduction")

    items = [
        "• IntelliFit is a comprehensive smart gym management platform that integrates "
        "AI-powered coaching with traditional gym operations.",
        "",
        "• The system provides personalized workout plan generation using fine-tuned "
        "machine learning models (Flan-T5 + LoRA).",
        "",
        "• It features a multi-role dashboard supporting Members, Coaches, "
        "Admins, and Receptionists.",
        "",
        "• Key capabilities include real-time equipment booking, InBody composition "
        "tracking, AI voice coaching, and a token-based payment economy.",
        "",
        "• The platform bridges the gap between AI-driven fitness recommendations "
        "and professional human coaching through a quality-control approval workflow.",
    ]
    add_bullet_content(slide, Inches(0.5), Inches(1.4), Inches(9), Inches(5.2),
                       items, font_size=15)
    add_footer(slide, 3, TOTAL_SLIDES)


def slide_04_problem(prs):
    """Slide 4: Problem Definition."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Problem Definition")

    add_section_header(slide, Inches(0.5), Inches(1.4), Inches(9),
                       "What Problem Does the Project Solve?")
    items1 = [
        "• Gym members often lack access to personalized, science-based workout plans "
        "tailored to their fitness level, goals, and available equipment.",
        "• Traditional gym management relies on manual processes for scheduling, "
        "equipment booking, and member tracking — leading to inefficiency.",
        "• Coach availability is limited; members may wait days for a customized program.",
    ]
    add_bullet_content(slide, Inches(0.5), Inches(1.9), Inches(9), Inches(2.2),
                       items1, font_size=14)

    add_section_header(slide, Inches(0.5), Inches(4.2), Inches(9),
                       "Why Is It Important?")
    items2 = [
        "• The global fitness industry is valued at $96+ billion (2024), "
        "with growing demand for digital fitness solutions.",
        "• 73% of gym members report higher retention when receiving personalized programs.",
        "• AI-powered fitness coaching can reduce trainer workload by up to 60%, "
        "while maintaining quality through human review.",
    ]
    add_bullet_content(slide, Inches(0.5), Inches(4.7), Inches(9), Inches(2),
                       items2, font_size=14)
    add_footer(slide, 4, TOTAL_SLIDES)


def slide_05_motivation(prs):
    """Slide 5: Motivation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Motivation")

    add_section_header(slide, Inches(0.5), Inches(1.4), Inches(4),
                       "Why We Chose This Project")
    items1 = [
        "• Passion for combining AI/ML with real-world health & fitness applications",
        "• Growing market need for intelligent gym management solutions",
        "• Opportunity to apply full-stack, ML, and cloud technologies in one project",
    ]
    add_bullet_content(slide, Inches(0.5), Inches(1.9), Inches(4.3), Inches(2.5),
                       items1, font_size=13)

    add_section_header(slide, Inches(5.2), Inches(1.4), Inches(4.5),
                       "Market & Academic Need")
    items2 = [
        "• Digital fitness market expected to reach $30B by 2026",
        "• Limited existing solutions that combine AI workout generation "
        "with coach quality control",
        "• Academic interest in practical NLP fine-tuning and MLOps pipelines",
    ]
    add_bullet_content(slide, Inches(5.2), Inches(1.9), Inches(4.5), Inches(2.5),
                       items2, font_size=13)

    add_section_header(slide, Inches(0.5), Inches(4.5), Inches(9),
                       "Key Benefits")
    items3 = [
        "• Personalized AI workout plans in under 2 seconds with 95%+ JSON validity",
        "• Reduced coach workload through AI-assisted program generation",
        "• Unified platform for all gym operations (booking, payments, tracking)",
        "• Real-time notifications and collaboration via SignalR",
    ]
    add_bullet_content(slide, Inches(0.5), Inches(5.0), Inches(9), Inches(2),
                       items3, font_size=13)
    add_footer(slide, 5, TOTAL_SLIDES)


def slide_06_objectives(prs):
    """Slide 6: Objectives."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Objectives")

    add_section_header(slide, Inches(0.5), Inches(1.4), Inches(4.5),
                       "Main Goals")
    goals = [
        "• Develop an AI-powered workout plan generator using fine-tuned NLP models",
        "• Build a multi-role gym management platform (Member, Coach, Admin, Reception)",
        "• Implement real-time equipment booking and availability tracking",
        "• Create a coach quality-control approval workflow for AI-generated plans",
        "• Design a token-based economy for gym services and payments",
    ]
    add_bullet_content(slide, Inches(0.5), Inches(1.9), Inches(4.3), Inches(4.5),
                       goals, font_size=13)

    add_section_header(slide, Inches(5.2), Inches(1.4), Inches(4.5),
                       "Key Features")
    features = [
        "• AI Workout Generator (Flan-T5 + LoRA fine-tuning)",
        "• AI Voice Coach (Google Gemini + Vapi)",
        "• InBody Composition Tracking with trend analysis",
        "• Real-time Chat & Notifications (SignalR)",
        "• Equipment Booking with time-slot management",
        "• Nutrition Plan Management",
        "• Progress Tracking & Achievements",
        "• Multi-role Dashboards with analytics",
    ]
    add_bullet_content(slide, Inches(5.2), Inches(1.9), Inches(4.5), Inches(4.5),
                       features, font_size=13)
    add_footer(slide, 6, TOTAL_SLIDES)


def slide_07_tech_comparison(prs):
    """Slide 7: Technology Comparison – Traditional vs AI-Powered Gym Management."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Technology Comparison: Traditional vs AI-Powered")

    # Table header
    headers = ["Aspect", "Traditional Approach", "IntelliFit (AI-Powered)"]
    rows = [
        ["Workout Plans", "Generic templates, manual creation",
         "AI-generated, personalized per user"],
        ["Coach Interaction", "In-person only, limited hours",
         "AI + human coach approval workflow"],
        ["Equipment Mgmt", "Paper-based / spreadsheets",
         "Real-time digital booking & tracking"],
        ["Body Tracking", "Manual measurements",
         "InBody integration with trend analytics"],
        ["Payments", "Cash / POS terminals",
         "Token economy with digital wallet"],
        ["Communication", "Phone / bulletin board",
         "Real-time SignalR notifications & chat"],
        ["Scalability", "Limited by staff count",
         "AI scales to unlimited members"],
    ]

    table = slide.shapes.add_table(
        len(rows) + 1, 3, Inches(0.4), Inches(1.5), Inches(9.2), Inches(5)
    ).table

    # Style header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(13)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY

    # Style data rows
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = val
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = DARK_TEXT
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY if r % 2 == 0 else WHITE

    # Set column widths
    table.columns[0].width = Inches(1.8)
    table.columns[1].width = Inches(3.5)
    table.columns[2].width = Inches(3.9)

    add_footer(slide, 7, TOTAL_SLIDES)


def slide_08_lifecycle(prs):
    """Slide 8: System Lifecycle / Workflow."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "System Lifecycle / Workflow")

    steps = [
        ("1", "Member Registration & Profile Setup",
         "User registers via Clerk OAuth, completes fitness profile"),
        ("2", "AI Workout Generation",
         "Member fills form → C# API → Python ML API → Flan-T5 generates plan"),
        ("3", "Coach Review & Approval",
         "Coach reviews AI-generated plan, approves or rejects with feedback"),
        ("4", "Plan Execution & Tracking",
         "Member follows plan, logs workouts, tracks progress & achievements"),
        ("5", "Equipment & Coach Booking",
         "Real-time booking system with time-slot management"),
        ("6", "InBody & Progress Monitoring",
         "Body composition tracking with trend analysis dashboards"),
        ("7", "Continuous AI Improvement",
         "Coach feedback fed back to ML model for retraining & improvement"),
    ]

    y_start = 1.4
    for i, (num, title, desc) in enumerate(steps):
        y = y_start + i * 0.78

        # Step number circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.4), Inches(y), Inches(0.45), Inches(0.45)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT
        circle.line.fill.background()
        tf = circle.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Title + description
        txBox = slide.shapes.add_textbox(Inches(1.1), Inches(y), Inches(8.5), Inches(0.7))
        tf2 = txBox.text_frame
        p1 = tf2.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = PRIMARY

        p2 = tf2.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = SUBTITLE_GRAY

    add_footer(slide, 8, TOTAL_SLIDES)


def slide_09_similar_apps(prs):
    """Slide 9: Similar Applications / Competitor Analysis."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Similar Applications")

    headers = ["Application", "AI Plans", "Coach Review", "Equipment Booking",
               "InBody", "Tokens"]
    rows = [
        ["MyFitnessPal", "✗", "✗", "✗", "✗", "✗"],
        ["Trainerize", "✗", "✓", "✗", "✗", "✗"],
        ["Mindbody", "✗", "✗", "✓", "✗", "✓"],
        ["Jefit", "Partial", "✗", "✗", "✗", "✗"],
        ["GymMaster", "✗", "✗", "✓", "✗", "✓"],
        ["IntelliFit (Ours)", "✓", "✓", "✓", "✓", "✓"],
    ]

    table = slide.shapes.add_table(
        len(rows) + 1, 6, Inches(0.3), Inches(1.5), Inches(9.4), Inches(4.5)
    ).table

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(12)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY

    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = val
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = DARK_TEXT
                paragraph.alignment = PP_ALIGN.CENTER
            is_last = r == len(rows) - 1
            cell.fill.solid()
            if is_last:
                cell.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
            elif r % 2 == 0:
                cell.fill.fore_color.rgb = LIGHT_GRAY
            else:
                cell.fill.fore_color.rgb = WHITE

    table.columns[0].width = Inches(2.2)
    for i in range(1, 6):
        table.columns[i].width = Inches(1.44)

    add_bullet_content(slide, Inches(0.5), Inches(6.2), Inches(9), Inches(0.7),
                       ["★ IntelliFit is the only platform combining AI workout generation, "
                        "coach quality control, and full gym management in one system."],
                       font_size=13, color=ACCENT_DARK)
    add_footer(slide, 9, TOTAL_SLIDES)


def slide_10_methodology(prs):
    """Slide 10: Methodology."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Methodology")

    add_section_header(slide, Inches(0.5), Inches(1.4), Inches(9),
                       "Development Methodology: Agile (Scrum)")
    items = [
        "• We adopted the Agile Scrum methodology with 2-week sprints.",
        "• Iterative development allowed continuous feedback integration "
        "from coaches and test users.",
        "• Sprint ceremonies: Daily stand-ups, Sprint Planning, Sprint Review, "
        "and Retrospectives.",
    ]
    add_bullet_content(slide, Inches(0.5), Inches(1.9), Inches(9), Inches(1.8),
                       items, font_size=14)

    add_section_header(slide, Inches(0.5), Inches(3.8), Inches(9),
                       "Why Agile Scrum?")
    reasons = [
        "• Rapid prototyping for AI model experimentation and iteration",
        "• Flexibility to pivot based on ML model performance and user testing results",
        "• Parallel development of frontend, backend, and ML components across team members",
        "• Continuous integration and deployment using Docker Compose",
        "• Incremental delivery of features with regular stakeholder demos",
    ]
    add_bullet_content(slide, Inches(0.5), Inches(4.3), Inches(9), Inches(2.5),
                       reasons, font_size=14)
    add_footer(slide, 10, TOTAL_SLIDES)


def slide_11_timeline(prs):
    """Slide 11: Timeline."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Project Timeline")

    phases = [
        ("Phase 1: Research & Planning", "Weeks 1–3",
         "Requirements gathering, technology selection, system design"),
        ("Phase 2: Database & Architecture", "Weeks 4–6",
         "PostgreSQL schema, Clean Architecture setup, API scaffolding"),
        ("Phase 3: Core Backend Development", "Weeks 7–12",
         "ASP.NET Core APIs, authentication, booking, payments"),
        ("Phase 4: ML Model Development", "Weeks 8–14",
         "Dataset creation, Flan-T5 fine-tuning, embeddings, API integration"),
        ("Phase 5: Frontend Development", "Weeks 10–16",
         "Next.js UI, dashboards, real-time features (SignalR)"),
        ("Phase 6: Integration & Testing", "Weeks 15–18",
         "End-to-end testing, Docker deployment, performance optimization"),
        ("Phase 7: Documentation & Presentation", "Weeks 19–20",
         "User documentation, API docs, presentation preparation"),
    ]

    y = 1.4
    for phase_title, duration, desc in phases:
        # Phase bar
        bar = add_shape_bg(slide, Inches(0.4), Inches(y), Inches(9.2),
                           Inches(0.65), LIGHT_GRAY)
        # Phase title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(y + 0.02),
                                         Inches(5), Inches(0.35))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = phase_title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = PRIMARY

        # Duration badge
        txBox2 = slide.shapes.add_textbox(Inches(7.5), Inches(y + 0.02),
                                          Inches(2), Inches(0.35))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = duration
        p2.font.size = Pt(12)
        p2.font.bold = True
        p2.font.color.rgb = ACCENT
        p2.alignment = PP_ALIGN.RIGHT

        # Description
        txBox3 = slide.shapes.add_textbox(Inches(0.7), Inches(y + 0.32),
                                          Inches(8.5), Inches(0.3))
        tf3 = txBox3.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = desc
        p3.font.size = Pt(11)
        p3.font.color.rgb = SUBTITLE_GRAY

        y += 0.78

    add_footer(slide, 11, TOTAL_SLIDES)


def slide_12_functional_req(prs):
    """Slide 12: Functional Requirements."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Functional Requirements")

    left_items = [
        "User Management",
        "  • Registration & OAuth (Clerk)",
        "  • Role-based access (4 roles)",
        "  • Profile management",
        "",
        "AI & Workout",
        "  • AI workout plan generation",
        "  • Coach review & approval workflow",
        "  • AI voice coaching (Gemini + Vapi)",
        "  • Workout logging & history",
        "",
        "Booking System",
        "  • Equipment time-slot booking",
        "  • Coach session scheduling",
        "  • Availability management",
    ]
    right_items = [
        "Health & Tracking",
        "  • InBody measurement tracking",
        "  • Progress analytics & trends",
        "  • Achievement & milestone system",
        "",
        "Payments & Tokens",
        "  • Token-based economy",
        "  • Package purchases",
        "  • Transaction history",
        "",
        "Communication",
        "  • Real-time notifications (SignalR)",
        "  • In-app chat system",
        "  • Activity feed",
        "",
        "Nutrition",
        "  • Meal & nutrition plan management",
    ]
    add_bullet_content(slide, Inches(0.3), Inches(1.4), Inches(4.7), Inches(5.5),
                       left_items, font_size=13, bold_first=True)
    add_bullet_content(slide, Inches(5.2), Inches(1.4), Inches(4.5), Inches(5.5),
                       right_items, font_size=13, bold_first=True)
    add_footer(slide, 12, TOTAL_SLIDES)


def slide_13_system_design(prs):
    """Slide 13: System Design."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "System Design")

    add_section_header(slide, Inches(0.5), Inches(1.4), Inches(9),
                       "System Architecture – Clean Architecture Pattern")
    arch_items = [
        "┌─────────────────────────────────────────────────────────┐",
        "│  Presentation Layer    │  Controllers, Hubs (SignalR)   │",
        "│  Service Layer         │  Business Logic, AI Services   │",
        "│  Domain Layer          │  Entities, Interfaces          │",
        "│  Infrastructure Layer  │  EF Core, PostgreSQL, Redis    │",
        "└─────────────────────────────────────────────────────────┘",
    ]
    add_bullet_content(slide, Inches(0.5), Inches(1.9), Inches(9), Inches(2.3),
                       arch_items, font_size=12, color=PRIMARY)

    add_section_header(slide, Inches(0.5), Inches(4.3), Inches(9),
                       "Microservices Overview")
    services = [
        "• Frontend Service (Next.js) ─────── Port 3000",
        "• Backend API (ASP.NET Core) ─────── Port 5000",
        "• ML Inference API (Python/Flask) ── Port 8000",
        "• TensorFlow Serving ──────────────── Port 8501",
        "• Embedding Service ───────────────── Port 5100",
        "• PostgreSQL + pgvector ────────────── Port 5432",
        "• Redis Cache ─────────────────────── Port 6379",
    ]
    add_bullet_content(slide, Inches(0.5), Inches(4.8), Inches(9), Inches(2.2),
                       services, font_size=12, color=DARK_TEXT)
    add_footer(slide, 13, TOTAL_SLIDES)


def slide_14_technologies(prs):
    """Slide 14: Technologies Used."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Technologies Used")

    categories = [
        ("Frontend", [
            "• Next.js 15 + React 19 + TypeScript",
            "• TailwindCSS + shadcn/ui",
            "• Clerk (Authentication & OAuth)",
            "• Convex (Real-time data sync)",
        ]),
        ("Backend", [
            "• ASP.NET Core (.NET 6/7) + C#",
            "• Entity Framework Core (ORM)",
            "• SignalR (Real-time communication)",
            "• JWT + OAuth 2.0 (Authorization)",
        ]),
        ("AI / Machine Learning", [
            "• Python + TensorFlow 2.15",
            "• Hugging Face Transformers (Flan-T5 + LoRA)",
            "• Sentence-Transformers (Embeddings)",
            "• Google Gemini + Vapi (Voice AI Coach)",
        ]),
        ("Infrastructure", [
            "• PostgreSQL 16 + pgvector extension",
            "• Docker + Docker Compose",
            "• Redis (Caching)",
            "• TensorFlow Serving (Model deployment)",
        ]),
    ]

    x_positions = [0.3, 5.1]
    y_positions = [1.4, 4.2]

    for idx, (cat_title, items) in enumerate(categories):
        col = idx % 2
        row = idx // 2
        x = x_positions[col]
        y = y_positions[row]

        # Category box
        box = add_shape_bg(slide, Inches(x), Inches(y), Inches(4.5),
                           Inches(2.4), LIGHT_GRAY)
        add_section_header(slide, Inches(x + 0.2), Inches(y + 0.1),
                           Inches(4), cat_title, font_size=16)
        add_bullet_content(slide, Inches(x + 0.2), Inches(y + 0.5),
                           Inches(4), Inches(1.8), items, font_size=12)
    add_footer(slide, 14, TOTAL_SLIDES)


def slide_15_database(prs):
    """Slide 15: Database Design."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Database Design")

    add_section_header(slide, Inches(0.5), Inches(1.4), Inches(9),
                       "PostgreSQL 16 with pgvector Extension")

    left_tables = [
        "Core Tables:",
        "  • Users (roles: Member, Coach, Admin, Reception)",
        "  • MemberProfiles / CoachProfiles",
        "  • WorkoutPlans + WorkoutPlanExercises",
        "  • WorkoutTemplates + WorkoutTemplateExercises",
        "  • Exercises (with 384-dim vector embeddings)",
        "  • Equipment + EquipmentCategories",
        "  • EquipmentTimeSlots",
        "  • Bookings",
    ]
    right_tables = [
        "Supporting Tables:",
        "  • InBodyMeasurements",
        "  • NutritionPlans + Meals + Ingredients",
        "  • TokenTransactions + TokenPackages",
        "  • Subscriptions + UserSubscriptions",
        "  • Payments",
        "  • Notifications + AuditLogs",
        "  • AiChatLogs + AiProgramGenerations",
        "  • UserMilestones + Achievements",
        "  • CoachReviews + WorkoutFeedback",
    ]
    add_bullet_content(slide, Inches(0.3), Inches(1.9), Inches(4.7), Inches(4),
                       left_tables, font_size=13, bold_first=True)
    add_bullet_content(slide, Inches(5.2), Inches(1.9), Inches(4.5), Inches(4),
                       right_tables, font_size=13, bold_first=True)

    key_features = [
        "Key Features: pgvector for exercise similarity search │ "
        "Seed data scripts for development │ EF Core migrations for schema management"
    ]
    add_bullet_content(slide, Inches(0.5), Inches(6.2), Inches(9), Inches(0.6),
                       key_features, font_size=12, color=ACCENT_DARK)
    add_footer(slide, 15, TOTAL_SLIDES)


def slide_16_ui_screens(prs):
    """Slide 16: UI / Screens."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "UI / Screens")

    roles = [
        ("Member Portal (15+ screens)", [
            "• Dashboard with fitness overview",
            "• AI Workout Generator form",
            "• Schedule & calendar view",
            "• InBody tracking with charts",
            "• Equipment booking interface",
            "• Token wallet & transactions",
            "• Progress & achievements",
        ]),
        ("Coach Portal (8+ screens)", [
            "• Coach dashboard with analytics",
            "• Client management",
            "• AI plan review & approval",
            "• Schedule management",
            "• Program creation tools",
        ]),
        ("Reception Portal (7+ screens)", [
            "• Member check-in system",
            "• New member registration",
            "• Booking management",
            "• Payment processing",
        ]),
        ("Admin Portal (5+ screens)", [
            "• Analytics dashboard",
            "• User & coach management",
            "• Equipment inventory",
            "• System-wide reports",
        ]),
    ]

    x_positions = [0.3, 5.1]
    y_positions = [1.4, 4.2]

    for idx, (role_title, items) in enumerate(roles):
        col = idx % 2
        row = idx // 2
        x = x_positions[col]
        y = y_positions[row]

        box = add_shape_bg(slide, Inches(x), Inches(y), Inches(4.5),
                           Inches(2.5), LIGHT_GRAY)
        add_section_header(slide, Inches(x + 0.2), Inches(y + 0.1),
                           Inches(4), role_title, font_size=14)
        add_bullet_content(slide, Inches(x + 0.2), Inches(y + 0.5),
                           Inches(4), Inches(1.9), items, font_size=11)
    add_footer(slide, 16, TOTAL_SLIDES)


def slide_17_pipeline(prs):
    """Slide 17: Project Pipeline / Workflow."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Project Pipeline / AI Workout Workflow")

    add_section_header(slide, Inches(0.5), Inches(1.4), Inches(9),
                       "AI Workout Generation Pipeline")
    pipeline = [
        "1. Member fills workout request form (days/week, fitness level, goals, equipment)",
        "2. Next.js frontend sends request to ASP.NET Core backend API",
        "3. Backend validates request and forwards to Python ML API (port 8000)",
        "4. Python API loads fine-tuned Flan-T5 + LoRA model",
        "5. Model generates personalized workout plan in structured JSON format",
        "6. Response validated (95%+ JSON validity) and returned to backend",
        "7. Backend stores plan in PostgreSQL and returns to frontend",
        "8. Plan displayed in interactive UI for member review",
    ]
    add_bullet_content(slide, Inches(0.5), Inches(1.9), Inches(9), Inches(2.8),
                       pipeline, font_size=13)

    add_section_header(slide, Inches(0.5), Inches(4.8), Inches(9),
                       "Coach Quality Control Loop")
    coach_flow = [
        "1. Coach receives notification of new AI-generated plan",
        "2. Reviews exercises, sets, reps, and progression logic",
        "3. Approves plan → Member can start workout immediately",
        "4. Rejects plan → Provides detailed feedback notes",
        "5. Feedback fed back to AI model for continuous improvement",
    ]
    add_bullet_content(slide, Inches(0.5), Inches(5.3), Inches(9), Inches(1.8),
                       coach_flow, font_size=13)
    add_footer(slide, 17, TOTAL_SLIDES)


def slide_18_progress(prs):
    """Slide 18: Progress Percentage."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Progress Percentage")

    modules = [
        ("Backend API (ASP.NET Core)", 95, "27 controllers, 30+ endpoints"),
        ("Frontend UI (Next.js)", 90, "42+ pages, 4 role dashboards"),
        ("AI/ML Models (Python)", 90, "Flan-T5 fine-tuned, embeddings ready"),
        ("Database Design", 95, "30+ tables, migrations, seed data"),
        ("Real-time Features (SignalR)", 85, "Notifications, chat"),
        ("Docker Deployment", 90, "6 containerized services"),
        ("Testing & QA", 80, "Unit, integration, API tests"),
        ("Documentation", 85, "API docs, ML guide, quick start"),
    ]

    y = 1.4
    for label, pct, detail in modules:
        # Label
        txBox = slide.shapes.add_textbox(Inches(0.4), Inches(y), Inches(3.5), Inches(0.3))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = PRIMARY

        # Progress bar background
        bar_bg = add_shape_bg(slide, Inches(4.2), Inches(y + 0.05),
                              Inches(4.5), Inches(0.25), LIGHT_GRAY)

        # Progress bar fill
        bar_width = 4.5 * (pct / 100)
        color = SUCCESS_GREEN if pct >= 90 else ACCENT
        bar_fill = add_shape_bg(slide, Inches(4.2), Inches(y + 0.05),
                                Inches(bar_width), Inches(0.25), color)

        # Percentage
        txBox2 = slide.shapes.add_textbox(Inches(8.8), Inches(y), Inches(0.7), Inches(0.3))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = f"{pct}%"
        p2.font.size = Pt(12)
        p2.font.bold = True
        p2.font.color.rgb = color

        # Detail text
        txBox3 = slide.shapes.add_textbox(Inches(4.2), Inches(y + 0.3),
                                          Inches(5), Inches(0.2))
        tf3 = txBox3.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = detail
        p3.font.size = Pt(9)
        p3.font.color.rgb = SUBTITLE_GRAY

        y += 0.67

    # Overall
    txBox = slide.shapes.add_textbox(Inches(3), Inches(6.6), Inches(4), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Overall Project Completion: ~89%"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = SUCCESS_GREEN
    p.alignment = PP_ALIGN.CENTER

    add_footer(slide, 18, TOTAL_SLIDES)


def slide_19_conclusion(prs):
    """Slide 19: Conclusion."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Conclusion")

    add_section_header(slide, Inches(0.5), Inches(1.4), Inches(9), "Results")
    results = [
        "• Successfully built a full-stack smart gym management system with AI capabilities",
        "• AI workout generator achieves 95%+ JSON validity with <2s latency",
        "• Multi-role platform serving Members, Coaches, Admins, and Receptionists",
        "• Real-time features powered by SignalR for notifications and chat",
        "• Containerized deployment with Docker Compose (6 services)",
    ]
    add_bullet_content(slide, Inches(0.5), Inches(1.9), Inches(9), Inches(2.3),
                       results, font_size=14)

    add_section_header(slide, Inches(0.5), Inches(4.3), Inches(9), "Achievements")
    achievements = [
        "• Fine-tuned Flan-T5 language model with LoRA (only 0.5% trainable parameters)",
        "• Implemented coach quality-control workflow — unique among competitor platforms",
        "• Built token-based economy for unified gym service payments",
        "• Designed 42+ frontend pages with responsive, modern UI",
        "• Created comprehensive API with 27 controllers and 30+ endpoints",
    ]
    add_bullet_content(slide, Inches(0.5), Inches(4.8), Inches(9), Inches(2),
                       achievements, font_size=14)
    add_footer(slide, 19, TOTAL_SLIDES)


def slide_20_challenges(prs):
    """Slide 20: Challenges."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Challenges & Solutions")

    headers = ["Challenge", "Solution"]
    rows = [
        ["ML model producing invalid JSON output",
         "Custom hybrid JSON parser + structured output formatting → 95%+ validity"],
        ["Large model size for deployment",
         "LoRA fine-tuning (0.5% params) + TensorFlow Serving for optimized inference"],
        ["Real-time sync across multiple clients",
         "SignalR WebSocket hubs for notifications, chat, and live updates"],
        ["Multi-role access control complexity",
         "Clerk OAuth + JWT with role-based middleware and separate dashboards"],
        ["Database performance with vector search",
         "PostgreSQL pgvector extension + indexed 384-dim embeddings"],
        ["Coordinating frontend, backend, and ML teams",
         "Docker Compose for unified dev environment + Agile sprints"],
    ]

    table = slide.shapes.add_table(
        len(rows) + 1, 2, Inches(0.3), Inches(1.5), Inches(9.4), Inches(5)
    ).table

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY

    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = val
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = DARK_TEXT
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY if r % 2 == 0 else WHITE

    table.columns[0].width = Inches(4.0)
    table.columns[1].width = Inches(5.4)

    add_footer(slide, 20, TOTAL_SLIDES)


def slide_21_future_work(prs):
    """Slide 21: Future Work."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Future Work")

    items = [
        "• Mobile Application – Native iOS/Android app using React Native or Flutter",
        "",
        "• Advanced AI Models – GPT-4 integration for more natural workout explanations "
        "and injury-aware plan modifications",
        "",
        "• Wearable Integration – Sync with smartwatches (Apple Watch, Fitbit) "
        "for real-time heart rate and calorie tracking during workouts",
        "",
        "• Computer Vision – Pose estimation for exercise form correction "
        "using device cameras",
        "",
        "• Social Features – Community challenges, leaderboards, "
        "and workout sharing among members",
        "",
        "• Multi-Gym Support – SaaS platform supporting multiple gym branches "
        "with centralized management",
        "",
        "• Advanced Analytics – Predictive models for member churn, "
        "optimal workout scheduling, and equipment utilization forecasting",
    ]
    add_bullet_content(slide, Inches(0.5), Inches(1.4), Inches(9), Inches(5.5),
                       items, font_size=14)
    add_footer(slide, 21, TOTAL_SLIDES)


def slide_22_thank_you(prs):
    """Slide 22: Thank You."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PRIMARY)

    # Thank you text
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You!"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(0.6))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "Questions & Discussion"
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(0xAA, 0xCC, 0xFF)
    p2.alignment = PP_ALIGN.CENTER

    # Accent line
    add_shape_bg(slide, Inches(3.5), Inches(4.0), Inches(3), Inches(0.05), ACCENT)

    # Contact info
    contact_items = [
        "IntelliFit – Smart Gym Management System",
        "AI-Powered Coaching • Equipment Booking • InBody Tracking",
        "",
        f"GitHub: {PROJECT_INFO['github_url']}",
    ]
    add_bullet_content(slide, Inches(1.5), Inches(4.5), Inches(7), Inches(2.5),
                       contact_items, font_size=14,
                       color=RGBColor(0xBB, 0xBB, 0xCC))


# ── Main ─────────────────────────────────────────────────────────────────────

def generate_presentation(output_path=None):
    """Generate the full IntelliFit graduation project presentation."""
    if output_path is None:
        output_path = os.path.join(
            os.path.dirname(os.path.abspath(__file__)),
            "..", "IntelliFit_Graduation_Presentation.pptx"
        )
    output_path = os.path.abspath(output_path)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_builders = [
        slide_01_title,
        slide_02_toc,
        slide_03_introduction,
        slide_04_problem,
        slide_05_motivation,
        slide_06_objectives,
        slide_07_tech_comparison,
        slide_08_lifecycle,
        slide_09_similar_apps,
        slide_10_methodology,
        slide_11_timeline,
        slide_12_functional_req,
        slide_13_system_design,
        slide_14_technologies,
        slide_15_database,
        slide_16_ui_screens,
        slide_17_pipeline,
        slide_18_progress,
        slide_19_conclusion,
        slide_20_challenges,
        slide_21_future_work,
        slide_22_thank_you,
    ]

    for builder in slide_builders:
        builder(prs)

    prs.save(output_path)
    print(f"Presentation generated successfully: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
    return output_path


if __name__ == "__main__":
    path = sys.argv[1] if len(sys.argv) > 1 else None
    generate_presentation(path)
