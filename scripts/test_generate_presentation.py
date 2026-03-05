"""
Tests for the IntelliFit graduation presentation generator.

Usage:
    python scripts/test_generate_presentation.py
"""

import os
import sys
import tempfile
import unittest

# Ensure scripts directory is on the path
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from generate_presentation import generate_presentation
from pptx import Presentation


class TestGeneratePresentation(unittest.TestCase):
    """Tests for the PowerPoint presentation generator."""

    @classmethod
    def setUpClass(cls):
        """Generate presentation once for all tests."""
        cls.output_path = os.path.join(tempfile.gettempdir(),
                                       "test_presentation.pptx")
        generate_presentation(cls.output_path)
        cls.prs = Presentation(cls.output_path)

    @classmethod
    def tearDownClass(cls):
        """Clean up generated file."""
        if os.path.exists(cls.output_path):
            os.remove(cls.output_path)

    def test_file_created(self):
        """Test that the .pptx file is created successfully."""
        self.assertTrue(os.path.exists(self.output_path))
        self.assertGreater(os.path.getsize(self.output_path), 0)

    def test_total_slides(self):
        """Test that exactly 22 slides are generated."""
        self.assertEqual(len(self.prs.slides), 22)

    def test_slide_dimensions(self):
        """Test that slide dimensions are 10x7.5 inches (widescreen)."""
        from pptx.util import Inches
        self.assertEqual(self.prs.slide_width, Inches(10))
        self.assertEqual(self.prs.slide_height, Inches(7.5))

    def test_title_slide_content(self):
        """Test that slide 1 contains the project title."""
        slide = self.prs.slides[0]
        all_text = self._get_slide_text(slide)
        self.assertIn("IntelliFit", all_text)
        self.assertIn("Smart Gym Management System", all_text)

    def test_table_of_contents(self):
        """Test that slide 2 is the Table of Contents."""
        slide = self.prs.slides[1]
        all_text = self._get_slide_text(slide)
        self.assertIn("Table of Contents", all_text)
        self.assertIn("Introduction", all_text)
        self.assertIn("Conclusion", all_text)

    def test_introduction_slide(self):
        """Test that the introduction slide has relevant content."""
        slide = self.prs.slides[2]
        all_text = self._get_slide_text(slide)
        self.assertIn("Introduction", all_text)
        self.assertIn("IntelliFit", all_text)

    def test_problem_definition_slide(self):
        """Test that the problem definition slide exists."""
        slide = self.prs.slides[3]
        all_text = self._get_slide_text(slide)
        self.assertIn("Problem Definition", all_text)

    def test_motivation_slide(self):
        """Test that the motivation slide exists."""
        slide = self.prs.slides[4]
        all_text = self._get_slide_text(slide)
        self.assertIn("Motivation", all_text)

    def test_objectives_slide(self):
        """Test that the objectives slide exists."""
        slide = self.prs.slides[5]
        all_text = self._get_slide_text(slide)
        self.assertIn("Objectives", all_text)

    def test_technology_comparison_slide(self):
        """Test that the tech comparison slide has a table."""
        slide = self.prs.slides[6]
        all_text = self._get_slide_text(slide)
        self.assertIn("Technology Comparison", all_text)
        # Verify table exists
        has_table = any(shape.has_table for shape in slide.shapes)
        self.assertTrue(has_table)

    def test_system_lifecycle_slide(self):
        """Test that the lifecycle slide exists."""
        slide = self.prs.slides[7]
        all_text = self._get_slide_text(slide)
        self.assertIn("System Lifecycle", all_text)

    def test_similar_applications_slide(self):
        """Test that the competitor analysis slide has a table."""
        slide = self.prs.slides[8]
        all_text = self._get_slide_text(slide)
        self.assertIn("Similar Applications", all_text)
        has_table = any(shape.has_table for shape in slide.shapes)
        self.assertTrue(has_table)

    def test_methodology_slide(self):
        """Test that the methodology slide mentions Agile."""
        slide = self.prs.slides[9]
        all_text = self._get_slide_text(slide)
        self.assertIn("Methodology", all_text)
        self.assertIn("Agile", all_text)

    def test_timeline_slide(self):
        """Test that the timeline slide has phases."""
        slide = self.prs.slides[10]
        all_text = self._get_slide_text(slide)
        self.assertIn("Timeline", all_text)
        self.assertIn("Phase 1", all_text)

    def test_functional_requirements_slide(self):
        """Test that functional requirements are listed."""
        slide = self.prs.slides[11]
        all_text = self._get_slide_text(slide)
        self.assertIn("Functional Requirements", all_text)

    def test_system_design_slide(self):
        """Test that system design mentions architecture."""
        slide = self.prs.slides[12]
        all_text = self._get_slide_text(slide)
        self.assertIn("System Design", all_text)
        self.assertIn("Clean Architecture", all_text)

    def test_technologies_slide(self):
        """Test that technologies slide lists key tech."""
        slide = self.prs.slides[13]
        all_text = self._get_slide_text(slide)
        self.assertIn("Technologies Used", all_text)
        self.assertIn("Next.js", all_text)
        self.assertIn("ASP.NET", all_text)
        self.assertIn("TensorFlow", all_text)

    def test_database_design_slide(self):
        """Test that database design mentions PostgreSQL."""
        slide = self.prs.slides[14]
        all_text = self._get_slide_text(slide)
        self.assertIn("Database Design", all_text)
        self.assertIn("PostgreSQL", all_text)

    def test_ui_screens_slide(self):
        """Test that UI slide lists multiple roles."""
        slide = self.prs.slides[15]
        all_text = self._get_slide_text(slide)
        self.assertIn("UI / Screens", all_text)
        self.assertIn("Member", all_text)
        self.assertIn("Coach", all_text)

    def test_pipeline_slide(self):
        """Test that the pipeline slide exists."""
        slide = self.prs.slides[16]
        all_text = self._get_slide_text(slide)
        self.assertIn("Pipeline", all_text)

    def test_progress_slide(self):
        """Test that the progress slide has percentages."""
        slide = self.prs.slides[17]
        all_text = self._get_slide_text(slide)
        self.assertIn("Progress", all_text)
        self.assertIn("%", all_text)

    def test_conclusion_slide(self):
        """Test that the conclusion slide exists."""
        slide = self.prs.slides[18]
        all_text = self._get_slide_text(slide)
        self.assertIn("Conclusion", all_text)

    def test_challenges_slide(self):
        """Test that the challenges slide has a table."""
        slide = self.prs.slides[19]
        all_text = self._get_slide_text(slide)
        self.assertIn("Challenges", all_text)
        has_table = any(shape.has_table for shape in slide.shapes)
        self.assertTrue(has_table)

    def test_future_work_slide(self):
        """Test that the future work slide exists."""
        slide = self.prs.slides[20]
        all_text = self._get_slide_text(slide)
        self.assertIn("Future Work", all_text)

    def test_thank_you_slide(self):
        """Test that the last slide is Thank You."""
        slide = self.prs.slides[21]
        all_text = self._get_slide_text(slide)
        self.assertIn("Thank You", all_text)
        self.assertIn("github.com/YoussefEssam74", all_text)

    def test_custom_output_path(self):
        """Test that a custom output path works."""
        custom_path = os.path.join(tempfile.gettempdir(),
                                   "custom_test.pptx")
        result = generate_presentation(custom_path)
        self.assertEqual(result, custom_path)
        self.assertTrue(os.path.exists(custom_path))
        os.remove(custom_path)

    def _get_slide_text(self, slide):
        """Extract all text from a slide."""
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    texts.append(p.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        texts.append(cell.text)
        return " ".join(texts)


if __name__ == "__main__":
    unittest.main()
